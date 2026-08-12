import os
import re
from pathlib import Path
import chromadb
from chromadb.utils import embedding_functions
from google import genai
from google.genai import types
from google.genai.errors import ClientError
from config.settings import GEMINI_API_KEY

# Document Parsers
import pypdf
import docx
from pptx import Presentation
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup


def load_prompt(filename: str) -> str:
    prompt_path = Path(__file__).parent.parent / "prompts" / filename
    if prompt_path.exists():
        return prompt_path.read_text(encoding="utf-8").strip()
    # Fallback default prompt if prompt file doesn't exist yet
    return "You are an expert document assistant. Answer the user's question accurately using only the provided document context chunks."


class DocumentAgent:
    def __init__(self, api_key: str = None):
        self.api_key = api_key or GEMINI_API_KEY
        if not self.api_key:
            raise ValueError("GEMINI_API_KEY is missing! Check your .env file.")

        self.client = genai.Client(api_key=self.api_key)
        self.model = "gemini-3.1-flash-lite"

        self.qa_prompt = load_prompt("doc_qa.txt")

        # Initialize ChromaDB vector store for Documents
        db_path = str(Path(__file__).parent.parent / "vectorstore" / "doc_db")
        self.chroma_client = chromadb.PersistentClient(path=db_path)
        self.embedding_fn = embedding_functions.DefaultEmbeddingFunction()

        self.current_doc_id = None
        self.current_collection = None

    def _clean_doc_id(self, file_path: str) -> str:
        """Generates a clean collection name ID from file path."""
        stem = Path(file_path).stem
        cleaned = re.sub(r'[^a-zA-Z0-9_-]', '_', stem)
        return cleaned[:50]  # Limit length for Chroma collection naming

    def extract_text_from_file(self, file_path: Path) -> tuple[str, str]:
        """Extracts text content based on file extension."""
        ext = file_path.suffix.lower()
        title = file_path.name
        text = ""

        if ext in [".txt", ".md", ".markdown"]:
            text = file_path.read_text(encoding="utf-8", errors="ignore")

        elif ext == ".pdf":
            reader = pypdf.PdfReader(file_path)
            pages = [page.extract_text() or "" for page in reader.pages]
            text = "\n".join(pages)

        elif ext == ".docx":
            doc = docx.Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

        elif ext == ".pptx":
            prs = Presentation(file_path)
            slide_texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
            text = "\n".join(slide_texts)

        elif ext == ".epub":
            book = epub.read_epub(str(file_path))
            chapters = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                soup = BeautifulSoup(item.get_content(), "html.parser")
                chapters.append(soup.get_text())
            text = "\n".join(chapters)

        elif ext in [".html", ".htm"]:
            soup = BeautifulSoup(file_path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
            text = soup.get_text()

        else:
            raise ValueError(f"Unsupported file format: {ext}")

        return title, text.strip()

    def fetch_and_index_document(self, raw_path: str) -> str:
        """Parses document, chunks text into windows, and embeds into ChromaDB."""
        file_path = Path(raw_path).expanduser().resolve()
        if not file_path.exists():
            return f"[ERROR] File not found: {raw_path}"

        doc_id = self._clean_doc_id(str(file_path))
        collection_name = f"doc_{doc_id}"

        # Check if document is already indexed in ChromaDB
        try:
            self.current_collection = self.chroma_client.get_collection(
                name=collection_name,
                embedding_function=self.embedding_fn
            )
            self.current_doc_id = doc_id
            return "ALREADY_INDEXED"
        except Exception:
            pass  # New document, process and index

        try:
            doc_title, full_text = self.extract_text_from_file(file_path)
        except Exception as e:
            return f"[ERROR] Failed to read document: {str(e)}"

        if not full_text:
            return "[ERROR] No readable text found in document."

        # Chunk text into ~300 word windows
        words = full_text.split()
        chunk_size = 300
        overlap = 50

        documents = []
        metadatas = []
        ids = []

        chunk_idx = 0
        for i in range(0, len(words), chunk_size - overlap):
            chunk_words = words[i:i + chunk_size]
            chunk_text = " ".join(chunk_words)

            if len(chunk_words) < 10:
                continue

            documents.append(chunk_text)
            metadatas.append({
                "doc_title": doc_title,
                "file_path": str(file_path),
                "chunk_index": chunk_idx
            })
            ids.append(f"chunk_{chunk_idx}")
            chunk_idx += 1

        if not documents:
            return "[ERROR] Could not generate valid text chunks."

        # Save to ChromaDB with metadata
        self.current_collection = self.chroma_client.create_collection(
            name=collection_name,
            embedding_function=self.embedding_fn,
            metadata={"title": doc_title, "file_path": str(file_path), "doc_id": doc_id}
        )
        self.current_collection.add(
            documents=documents,
            metadatas=metadatas,
            ids=ids
        )
        self.current_doc_id = doc_id
        return "INDEX_SUCCESS"

    def list_indexed_docs(self) -> list[dict]:
        """Lists all document collections stored in ChromaDB."""
        collections = self.chroma_client.list_collections()
        doc_list = []
        for col in collections:
            if col.name.startswith("doc_"):
                doc_id = col.name.replace("doc_", "")
                title = f"Document ({doc_id})"
                file_path = ""

                if col.metadata:
                    title = col.metadata.get("title", title)
                    file_path = col.metadata.get("file_path", "")

                doc_list.append({
                    "doc_id": doc_id,
                    "title": title,
                    "file_path": file_path,
                    "chunk_count": col.count()
                })
        return doc_list

    def switch_active_doc(self, target: str) -> str:
        """Switches active RAG session to an existing document in vector store."""
        target_clean = self._clean_doc_id(target)
        collection_name = f"doc_{target_clean}"

        existing_cols = [c.name for c in self.chroma_client.list_collections()]
        if collection_name in existing_cols:
            self.current_collection = self.chroma_client.get_collection(
                name=collection_name,
                embedding_function=self.embedding_fn
            )
            self.current_doc_id = target_clean
            return f"[SUCCESS] Switched active session to Document [{target_clean}]."
        else:
            return f"[ERROR] Document [{target}] is not indexed. Load it using: /doc <file_path>"

    def auto_load_latest(self) -> str | None:
        """Restores the most recently indexed document session on CLI start."""
        indexed = self.list_indexed_docs()
        if indexed:
            latest_id = indexed[-1]["doc_id"]
            self.switch_active_doc(latest_id)
            return latest_id
        return None

    def query_rag(self, query: str, top_k: int = 4) -> str:
        """Retrieves relevant context chunks from ChromaDB for user prompt."""
        if not self.current_collection:
            return "[ERROR] No active document loaded."

        results = self.current_collection.query(
            query_texts=[query],
            n_results=top_k
        )
        retrieved_chunks = results["documents"][0]
        return "\n\n---\n\n".join(retrieved_chunks)

    def process_doc_query(self, file_or_input: str, user_question: str = None) -> str:
        """Processes document loading, summaries, or Q&A."""
        # 1. Direct Question on active document
        if self.current_collection and user_question:
            context = self.query_rag(query=user_question, top_k=4)
            contents = f"User Question: {user_question}\n\nDocument Context:\n{context}"
        
        # 2. File path passed for indexing or initial summary
        else:
            status = self.fetch_and_index_document(file_or_input)
            if status.startswith("[ERROR]"):
                return status

            if not user_question:
                context = self.query_rag(query="executive summary main topics overview key points", top_k=6)
                contents = f"Provide a complete, structured summary of this document based on these context passages:\n\n{context}"
            else:
                context = self.query_rag(query=user_question, top_k=4)
                contents = f"User Question: {user_question}\n\nDocument Context:\n{context}"

        try:
            response = self.client.models.generate_content(
                model=self.model,
                contents=contents,
                config=types.GenerateContentConfig(
                    system_instruction=self.qa_prompt,
                    temperature=0.1
                )
            )
            return response.text
        except ClientError as e:
            if e.code == 429 or "RESOURCE_EXHAUSTED" in str(e):
                return "[ERROR] Gemini API Quota Exceeded (429 Rate Limit)."
            return f"[ERROR] API Client Error: {str(e)}"
        except Exception as e:
            return f"[ERROR] Generation failed: {str(e)}"